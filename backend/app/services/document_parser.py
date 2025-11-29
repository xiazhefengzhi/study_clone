"""
Document parsing service for extracting text content from PDF, PPT, and Word files
"""
import io
from typing import Optional
from pathlib import Path

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    from pptx import Presentation
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

try:
    import docx
    HAS_PYTHON_DOCX = True
except ImportError:
    HAS_PYTHON_DOCX = False


class DocumentParser:
    """Parse various document formats to extract text content"""

    @staticmethod
    def parse_pdf(file_content: bytes) -> str:
        """
        Extract text from PDF file

        Args:
            file_content: PDF file content as bytes

        Returns:
            Extracted text content
        """
        if not HAS_PDFPLUMBER:
            raise ImportError("pdfplumber is not installed. Install with: pip install pdfplumber")

        text_content = []

        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_content.append(text)

        return "\n\n".join(text_content)

    @staticmethod
    def parse_pptx(file_content: bytes) -> str:
        """
        Extract text from PowerPoint file

        Args:
            file_content: PPTX file content as bytes

        Returns:
            Extracted text content
        """
        if not HAS_PYTHON_PPTX:
            raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")

        text_content = []
        prs = Presentation(io.BytesIO(file_content))

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"## Slide {slide_num}"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            text_content.append("\n".join(slide_text))

        return "\n\n".join(text_content)

    @staticmethod
    def parse_docx(file_content: bytes) -> str:
        """
        Extract text from Word document

        Args:
            file_content: DOCX file content as bytes

        Returns:
            Extracted text content
        """
        if not HAS_PYTHON_DOCX:
            raise ImportError("python-docx is not installed. Install with: pip install python-docx")

        doc = docx.Document(io.BytesIO(file_content))
        text_content = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_content.append(" | ".join(row_text))

        return "\n\n".join(text_content)

    @staticmethod
    def parse_document(file_content: bytes, file_extension: str) -> str:
        """
        Parse document based on file extension

        Args:
            file_content: File content as bytes
            file_extension: File extension (e.g., '.pdf', '.pptx', '.docx')

        Returns:
            Extracted text content

        Raises:
            ValueError: If file format is not supported
            ImportError: If required library is not installed
        """
        extension = file_extension.lower()

        if extension == ".pdf":
            return DocumentParser.parse_pdf(file_content)
        elif extension in [".ppt", ".pptx"]:
            return DocumentParser.parse_pptx(file_content)
        elif extension in [".doc", ".docx"]:
            return DocumentParser.parse_docx(file_content)
        else:
            raise ValueError(f"Unsupported file format: {extension}")

    @staticmethod
    async def parse_from_storage(storage_path: str) -> str:
        """
        Parse document from storage path

        Args:
            storage_path: Path to file in storage

        Returns:
            Extracted text content
        """
        from app.services.storage_service import storage_service

        # Download file from storage
        file_content = await storage_service.download_file(storage_path)

        # Get file extension
        file_extension = Path(storage_path).suffix

        # Parse document
        return DocumentParser.parse_document(file_content, file_extension)


# Singleton instance
document_parser = DocumentParser()
